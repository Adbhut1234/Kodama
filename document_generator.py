from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import os
import random
import urllib.request
from io import BytesIO

from utils import (
    ask_llm,
    search_internet,
    should_search,
    chat_with_search,
    sanitize_filename,
)

# ── Output Path ───────────────────────────────────────────────────────────────
OUTPUT_DIR = os.path.join(os.path.expanduser('~'), 'Documents', 'Kodama')
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

# ── Color Themes (randomly selected per presentation) ─────────────────────────
THEMES = [
    {   # Midnight Indigo
        "bg":      RGBColor(18,  18,  35),
        "accent":  RGBColor(99,  102, 241),
        "accent2": RGBColor(168, 85,  247),
        "white":   RGBColor(255, 255, 255),
        "light":   RGBColor(200, 200, 220),
        "card":    RGBColor(30,  30,  60),
        "card2":   RGBColor(45,  40,  80),
    },
    {   # Ocean Teal
        "bg":      RGBColor(10,  20,  30),
        "accent":  RGBColor(0,   200, 180),
        "accent2": RGBColor(0,   140, 200),
        "white":   RGBColor(255, 255, 255),
        "light":   RGBColor(180, 220, 215),
        "card":    RGBColor(18,  35,  45),
        "card2":   RGBColor(25,  50,  60),
    },
    {   # Ember Red
        "bg":      RGBColor(25,  12,  15),
        "accent":  RGBColor(239, 68,  68),
        "accent2": RGBColor(249, 115, 22),
        "white":   RGBColor(255, 255, 255),
        "light":   RGBColor(220, 200, 200),
        "card":    RGBColor(45,  22,  28),
        "card2":   RGBColor(60,  30,  35),
    },
    {   # Emerald Forest
        "bg":      RGBColor(10,  22,  18),
        "accent":  RGBColor(16,  185, 129),
        "accent2": RGBColor(52,  211, 153),
        "white":   RGBColor(255, 255, 255),
        "light":   RGBColor(200, 225, 215),
        "card":    RGBColor(20,  40,  32),
        "card2":   RGBColor(28,  55,  45),
    },
    {   # Royal Gold
        "bg":      RGBColor(20,  18,  12),
        "accent":  RGBColor(234, 179, 8),
        "accent2": RGBColor(245, 158, 11),
        "white":   RGBColor(255, 255, 255),
        "light":   RGBColor(225, 220, 200),
        "card":    RGBColor(38,  35,  22),
        "card2":   RGBColor(50,  45,  28),
    },
]

# ── Image Downloader ──────────────────────────────────────────────────────────
def download_topic_images(topic, count=4):
    """Download images and save to temp folder, returning file paths."""
    image_paths = []
    temp_dir = os.path.join(OUTPUT_DIR, "temp_images")
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
        
    try:
        from ddgs import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.images(topic, max_results=count + 2))
        
        for idx, r in enumerate(results[:count + 2]):
            if len(image_paths) >= count:
                break
            try:
                url = r.get('image', '')
                if not url: continue
                
                req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                data = urllib.request.urlopen(req, timeout=5).read()
                
                if len(data) > 10000:
                    ext = ".jpg" if "jpg" in url.lower() or "jpeg" in url.lower() else ".png"
                    path = os.path.join(temp_dir, f"img_{idx}{ext}")
                    with open(path, "wb") as f:
                        f.write(data)
                    image_paths.append(path)
            except: continue
    except Exception as e:
        print(f"⚠️ Image download skipped: {e}")
    return image_paths

# ── PPT Helper Functions ──────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color):
    """Add a filled rectangle to slide"""
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size, color, bold=False, italic=False, align=PP_ALIGN.LEFT):
    """Add a styled text box to slide"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Segoe UI"
    return txBox


# ── PPT Generator ─────────────────────────────────────────────────────────────
def create_ppt(topic, num_slides=5):
    print(f"\n📊 Processing PPT Synthesis ({num_slides} slides) for: {topic}")
    # Always search for deep context for documents
    search_context = ""
    print("🔍 Fetching deep context from the web...")
    results = search_internet(topic)
    if results:
        search_context = f"Use this detailed information to make the presentation highly informative:\n{results}"

    # Scale token budget — ~200 tokens per slide is a safe estimate
    token_budget = max(1500, num_slides * 200)

    # ── High-End Design Prompt ────────────────────────────────────────────────
    prompt = f"""
Create a premium, high-impact {num_slides} slide presentation about: {topic}
{search_context}

CRITICAL: Do NOT just list bullet points. Use a variety of data-rich layouts.
Available Layouts (YOU MUST STRICTLY USE THESE EXACT KEYS):
- "hero": Requires keys: "heading", "emoji", "quote"
- "metric": Requires keys: "heading", "metrics" (array of {"value", "label"})
- "process": Requires keys: "heading", "steps" (array of {"title", "desc"})
- "feature_grid": Requires keys: "heading", "features" (array of {"title", "desc"})
- "standard": Requires keys: "heading", "points" (array of strings)

Structure your JSON exactly like this:
{{
    "title": "Main Title",
    "subtitle": "Secondary Title",
    "design_style": "futuristic / minimal / executive",
    "slides": [
        {{
            "heading": "Heading",
            "layout": "metric",
            "emoji": "📊",
            "metrics": [
                {{"value": "98%", "label": "Accuracy Rating"}},
                {{"value": "2.4x", "label": "Performance Gain"}}
            ]
        }},
        {{
            "heading": "Heading",
            "layout": "process",
            "steps": [
                {{"title": "Initialize", "desc": "Context setup"}},
                {{"title": "Process", "desc": "Neural mapping"}}
            ]
        }}
    ]
}}
"""

    raw = ask_llm(prompt, temperature=0.4, num_predict=token_budget, format="json")

    try:
        import re
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        if match:
            json_str = match.group(0)
            data = json.loads(json_str)
        else:
            raise ValueError("No JSON block found")
    except Exception as e:
        print(f"JSON Parse Error: {e}")
        data = {"title": topic, "subtitle": "Synthesis Archive", "design_style": "futuristic", "slides": []}
        for i in range(num_slides):
            data["slides"].append({"heading": f"Insight {i+1}", "layout": "hero", "emoji": "🚀", "quote": "Innovation distinguishes between a leader and a follower."})

    # ── Build Presentation ────────────────────────────────────────────────────
    T = random.choice(THEMES)
    
    # Download images
    images = download_topic_images(topic, count=5)
    
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── TITLE SLIDE ───────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, T["bg"])
    if images:
        try:
            img = images[0]; img.seek(0)
            slide.shapes.add_picture(img, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            # Dark Gradient Overlay
            ov = add_rect(slide, 0, 0, 13.33, 7.5, T["bg"])
            from pptx.oxml.ns import qn
            ov.fill.solid()
            alpha = ov.fill._fill.find(qn('a:srgbClr')).makeelement(qn('a:alpha'), {'val': '65000'})
            ov.fill._fill.find(qn('a:srgbClr')).append(alpha)
        except: pass

    add_text(slide, data["title"], 1, 2.5, 11.33, 2, font_size=54, color=T["white"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, data["subtitle"], 1, 4.5, 11.33, 0.8, font_size=24, color=T["accent"], align=PP_ALIGN.CENTER)
    add_rect(slide, 5.16, 5.5, 3, 0.08, T["accent"])

    # ── CONTENT SLIDES ────────────────────────────────────────────────────────
    for idx, sdata in enumerate(data["slides"]):
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, T["bg"])
        
        layout = sdata.get("layout", "hero")
        heading = sdata.get("heading", "")
        emoji = sdata.get("emoji", "💡")
        img = images[(idx+1) % len(images)] if images else None

        # Side bar accent
        add_rect(slide, 0, 0, 0.08, 7.5, T["accent"])
        
        if layout == "hero":
            # ── HERO LAYOUT (Big impact) ──
            add_text(slide, emoji, 1, 1.5, 11.33, 2, font_size=96, color=T["white"], align=PP_ALIGN.CENTER)
            add_text(slide, heading, 1, 3.8, 11.33, 1, font_size=44, color=T["white"], bold=True, align=PP_ALIGN.CENTER)
            quote = sdata.get("quote", sdata.get("sentence", sdata.get("text", "")))
            add_text(slide, quote, 2, 5.0, 9.33, 1, font_size=22, color=T["light"], italic=True, align=PP_ALIGN.CENTER)

        elif layout == "metric":
            # ── METRIC LAYOUT (Stats) ──
            add_text(slide, heading, 1, 0.8, 11, 1, font_size=32, color=T["white"], bold=True)
            metrics = sdata.get("metrics", [])[:3]
            for i, m in enumerate(metrics):
                mx = 1.0 + i * 4.0
                add_rect(slide, mx, 2.5, 3.5, 3.5, T["card"])
                add_text(slide, m.get("value", "0"), mx, 3.2, 3.5, 1.2, font_size=54, color=T["accent"], bold=True, align=PP_ALIGN.CENTER)
                add_text(slide, m.get("label", ""), mx + 0.25, 4.8, 3, 0.8, font_size=16, color=T["light"], align=PP_ALIGN.CENTER)

        elif layout == "process":
            # ── PROCESS LAYOUT (Flow) ──
            add_text(slide, heading, 1, 0.8, 11, 1, font_size=32, color=T["white"], bold=True)
            steps = sdata.get("steps", [])[:4]
            # Connecting line
            add_rect(slide, 1.5, 3.8, 10, 0.04, T["light"])
            for i, step in enumerate(steps):
                sx = 1.5 + i * 2.8
                # Number circle
                circ = slide.shapes.add_shape(9, Inches(sx), Inches(3.5), Inches(0.6), Inches(0.6))
                circ.fill.solid(); circ.fill.fore_color.rgb = T["accent"]; circ.line.fill.background()
                add_text(slide, str(i+1), sx, 3.55, 0.6, 0.6, font_size=18, color=T["white"], bold=True, align=PP_ALIGN.CENTER)
                # Text
                add_text(slide, step.get("title", ""), sx - 0.5, 4.3, 1.6, 0.6, font_size=16, color=T["white"], bold=True, align=PP_ALIGN.CENTER)
                desc = step.get("desc", step.get("description", step.get("content", "")))
                add_text(slide, desc, sx - 0.8, 5.0, 2.2, 1.0, font_size=11, color=T["light"], align=PP_ALIGN.CENTER)

        elif layout == "comparison_pro":
            # ── COMPARISON PRO ──
            add_rect(slide, 0.5, 1.5, 6.0, 5.5, T["card"])
            add_rect(slide, 6.8, 1.5, 6.0, 5.5, T["card2"])
            data_c = sdata.get("comparison", {})
            opt_a = data_c.get("option_a", {"title": "Option A", "points": ["Point 1"]})
            opt_b = data_c.get("option_b", {"title": "Option B", "points": ["Point 1"]})
            
            add_text(slide, opt_a.get("title", "OPTION A"), 1.0, 2.0, 5, 0.8, font_size=24, color=T["accent"], bold=True)
            add_text(slide, opt_b.get("title", "OPTION B"), 7.3, 2.0, 5, 0.8, font_size=24, color=T["accent2"], bold=True)
            
            y_a = 3.0
            for p in opt_a.get("points", [])[:4]:
                add_text(slide, f"• {p}", 1.0, y_a, 5, 0.8, font_size=16, color=T["light"])
                y_a += 0.8
                
            y_b = 3.0
            for p in opt_b.get("points", [])[:4]:
                add_text(slide, f"• {p}", 7.3, y_b, 5, 0.8, font_size=16, color=T["light"])
                y_b += 0.8

        elif layout == "feature_grid":
            # ── FEATURE GRID (4 Cards) ──
            add_text(slide, heading, 1, 0.6, 11, 1, font_size=30, color=T["white"], bold=True, align=PP_ALIGN.CENTER)
            features = sdata.get("features", [])[:4]
            for i, f in enumerate(features):
                cx, cy = 1.2 + (i%2)*5.8, 2.0 + (i//2)*2.4
                add_rect(slide, cx, cy, 5.2, 2.0, T["card"])
                add_rect(slide, cx, cy, 0.1, 2.0, T["accent"] if i<2 else T["accent2"])
                add_text(slide, f.get("title", "Feature"), cx + 0.4, cy + 0.3, 4.5, 0.6, font_size=18, color=T["white"], bold=True)
                desc = f.get("desc", f.get("description", f.get("content", "")))
                add_text(slide, desc, cx + 0.4, cy + 0.9, 4.5, 1.0, font_size=12, color=T["light"])

        else:
            # ── FALLBACK / STANDARD (Full width image split) ──
            if img:
                try:
                    img.seek(0); slide.shapes.add_picture(img, Inches(0.08), Inches(0), Inches(5.5), Inches(7.5))
                    cx, cw = 6.0, 6.5
                except: cx, cw = 1.0, 11.3
            else: cx, cw = 1.0, 11.3
            
            add_text(slide, heading, cx, 1.2, cw, 1, font_size=36, color=T["white"], bold=True)
            add_rect(slide, cx, 2.4, 4, 0.06, T["accent"])
            
            # Smart fallback for standard bullets
            points = sdata.get("points", sdata.get("bullets", sdata.get("content", [])))
            if not isinstance(points, list):
                points = [str(points)]
            if not points:
                points = ["Generating sub-neural pathways...", "Awaiting data alignment..."]
                
            py = 3.0
            for p in points[:4]:
                add_text(slide, f"❯  {p}", cx, py, cw, 0.8, font_size=18, color=T["light"])
                py += 1.0

        # Footer slide number
        add_text(slide, f"{idx+1}", 12.5, 6.8, 0.5, 0.5, font_size=14, color=T["accent"], bold=True, align=PP_ALIGN.RIGHT)

    # ── Save ──────────────────────────────────────────────────────────────────
    filename = f"{sanitize_filename(topic)}.pptx"
    full_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(full_path)
    print(f"✅ Professional PPT saved: {full_path}")
    return full_path


# ── PDF Generator ─────────────────────────────────────────────────────────────
def create_pdf(topic):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor, white, slateblue, slategray
    from reportlab.lib.units import mm
    from reportlab.pdfgen import canvas

    print(f"\n💎 Redesigning PDF with Premium Executive Layout: {topic}")

    # ── Search & Synthesis ───────────────────────────────────────────────────
    search_results = search_internet(topic)
    search_context = f"USE THESE SEARCH RESULTS FOR DATA:\n{search_results}\n" if search_results else ""

    prompt = f"""
{search_context}
Generate a comprehensive technical report/notes about "{topic}".
Provide technical definitions, key concepts, and detailed explanations.

CRITICAL INSTRUCTION: Your response must be a JSON object inside a code block.
Even if you are a small model, try your best to follow this structure:
{{
    "title": "Clear Title",
    "subtitle": "Informative Subtitle",
    "abstract": "Executive summary or overview",
    "chapters": [
        {{"heading": "Concept 1", "body": "Technical content..."}},
        {{"heading": "Concept 2", "body": "Technical content..."}}
    ],
    "key_takeaway": "Summary sentence"
}}
"""
    try:
        from utils import clean_json_string
        print("🧠 Synthesizing neural data (Optimized for Universal Compatibility)...")
        raw_data = ask_llm(prompt, num_predict=3000, temperature=0.4)
        
        # Repair and Parse
        clean_json = clean_json_string(raw_data)
        data = json.loads(clean_json)
            
    except Exception as e:
        print(f"⚠️ Synthesis Error: {e}. Reverting to Direct Data Stream...")
        # Smart Fallback: If JSON fails, at least provide the raw text and search info
        content_snippet = raw_data if len(raw_data) > 100 else "The neural model failed to synthesize a response. Check your connection to Ollama."
        data = {
            "title": str(topic).upper(),
            "subtitle": "DIRECT DATA SYNTHESIS",
            "abstract": "Automatic synthesis initiated after structure failure. This usually happens with very small models (e.g. 1B-3B) or connection timeouts.",
            "chapters": [
                {"heading": "Direct Analysis", "body": str(content_snippet)},
                {"heading": "Internet Context", "body": str(search_results) if search_results else "No internet data available."}
            ],
            "key_takeaway": "Synthesis completed via direct stream."
        }

    # ── PDF Themes (Randomly selected for each report) ───────────────────────
    PDF_THEMES = [
        {"primary": "#4f46e5", "secondary": "#1e293b", "accent": "#94a3b8", "name": "Midnight Indigo"},
        {"primary": "#059669", "secondary": "#064e3b", "accent": "#6ee7b7", "name": "Emerald Executive"},
        {"primary": "#dc2626", "secondary": "#450a0a", "accent": "#f87171", "name": "Ember Whitepaper"},
        {"primary": "#0284c7", "secondary": "#082f49", "accent": "#7dd3fc", "name": "Oceanic Strategy"},
        {"primary": "#7c3aed", "secondary": "#2e1065", "accent": "#c4b5fd", "name": "Royal Synthesis"},
        {"primary": "#d97706", "secondary": "#451a03", "accent": "#fbbf24", "name": "Amber Intelligence"}
    ]
    theme = random.choice(PDF_THEMES)
    p_color = HexColor(theme["primary"])
    s_color = HexColor(theme["secondary"])
    a_color = HexColor(theme["accent"])
    
    print(f"🎨 Applying Theme: {theme['name']}")

    # ── PDF Setup ────────────────────────────────────────────────────────────
    filename = f"{sanitize_filename(topic)}.pdf"
    full_path = os.path.join(OUTPUT_DIR, filename)
    
    doc = SimpleDocTemplate(
        full_path,
        pagesize=A4,
        rightMargin=25*mm,
        leftMargin=25*mm,
        topMargin=30*mm,
        bottomMargin=30*mm
    )

    styles = getSampleStyleSheet()
    
    # Custom Typography
    style_h1 = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=20, textColor=HexColor('#1e293b'), spaceAfter=12, fontName='Helvetica-Bold')
    style_h2 = ParagraphStyle('H2', parent=styles['Heading2'], fontSize=14, textColor=p_color, spaceBefore=20, spaceAfter=8, fontName='Helvetica-Bold')
    style_body = ParagraphStyle('Body', parent=styles['Normal'], fontSize=11, textColor=HexColor('#475569'), leading=16, alignment=4, spaceAfter=10)
    style_quote = ParagraphStyle('Quote', parent=style_body, leftIndent=10, borderPadding=10, backColor=HexColor('#f8fafc'), textColor=p_color, fontName='Helvetica-Oblique')

    # ── ELEMENTS ─────────────────────────────────────────────────────────────
    elements = []

    # ── COVER PAGE (Empty, handled by draw_cover) ──
    elements.append(PageBreak())

    # ── CONTENT PAGE ──
    elements.append(Paragraph("EXECUTIVE SUMMARY", style_h2))
    elements.append(Paragraph(data.get("abstract", ""), style_body))
    elements.append(Spacer(1, 10*mm))

    for ch in data.get("chapters", []):
        elements.append(Paragraph(ch.get("heading", "").upper(), style_h2))
        elements.append(Paragraph(ch.get("body", ""), style_body))
    
    elements.append(Spacer(1, 10*mm))
    elements.append(Paragraph("STRATEGIC CONCLUSION", style_h2))
    elements.append(Paragraph(data.get("key_takeaway", ""), style_quote))

    # ── FIXED ELEMENTS (Cover, Header, Footer) ──
    def draw_all(canvas, doc):
        if doc.page == 1:
            # ── DESIGNER COVER ──
            canvas.saveState()
            canvas.setFillColor(p_color)
            canvas.rect(0, 0, 8*mm, A4[1], fill=1, stroke=0) 
            canvas.setFillColor(s_color)
            canvas.rect(8*mm, A4[1]-100*mm, A4[0], 100*mm, fill=1, stroke=0) 
            
            canvas.setFillColor(white)
            canvas.setFont('Helvetica-Bold', 32)
            title_text = str(data.get("title", topic) or topic).upper()
            canvas.drawString(25*mm, A4[1]-60*mm, title_text)
            
            canvas.setFont('Helvetica', 14)
            canvas.setFillColor(a_color)
            subtitle_text = str(data.get("subtitle", "INTELLIGENCE REPORT") or "INTELLIGENCE REPORT").upper()
            canvas.drawString(25*mm, A4[1]-75*mm, subtitle_text)
            
            canvas.setFillColor(s_color)
            canvas.setFont('Helvetica-Bold', 18)
            canvas.drawString(25*mm, 40*mm, "KODAMA AI")
            canvas.setFont('Helvetica', 10)
            canvas.setFillColor(HexColor('#64748b'))
            canvas.drawString(25*mm, 34*mm, f"{theme['name'].upper()} — SYNTHESIS ENGINE")
            
            canvas.setStrokeColor(p_color)
            canvas.setLineWidth(1)
            canvas.line(25*mm, 50*mm, 100*mm, 50*mm)
            canvas.restoreState()
        else:
            # ── STANDARD PAGES HEADER/FOOTER ──
            canvas.saveState()
            canvas.setFont('Helvetica-Bold', 8)
            canvas.setFillColor(HexColor('#94a3b8'))
            canvas.drawString(25*mm, A4[1]-15*mm, data.get("title", "").upper())
            canvas.drawRightString(A4[0]-25*mm, A4[1]-15*mm, f"{theme['name'].upper()} | CONFIDENTIAL")
            
            canvas.setStrokeColor(HexColor('#e2e8f0'))
            canvas.line(25*mm, 20*mm, A4[0]-25*mm, 20*mm)
            canvas.setFont('Helvetica', 8)
            canvas.drawString(25*mm, 15*mm, "Generated by Kodama AI Synthesis Engine")
            canvas.drawRightString(A4[0]-25*mm, 15*mm, f"Page {doc.page}")
            canvas.restoreState()

    # Build PDF
    doc.build(elements, onFirstPage=draw_all, onLaterPages=draw_all)
    print(f"✅ Executive Redesigned PDF saved: {full_path}")
    return full_path


# ── Main Loop ─────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 55)
    print("🚀  AI Assistant — Chat + Docs + Internet Search")
    print("=" * 55)
    print("Commands:")
    print("  'make ppt: topic'  →  Create PowerPoint file")
    print("  'make pdf: topic'  →  Create PDF report file")
    print("  'exit'             →  Quit")
    print("=" * 55 + "\n")

    while True:
        user_input = input("You: ").strip()

        if not user_input:
            continue

        if user_input.lower() == "exit":
            print("Goodbye! 👋")
            break

        elif user_input.lower().startswith("make ppt:"):
            topic = user_input[9:].strip()
            if topic:
                create_ppt(topic)
            else:
                print("⚠️ Please provide a topic. Example: make ppt: Artificial Intelligence")

        elif user_input.lower().startswith("make pdf:"):
            topic = user_input[9:].strip()
            if topic:
                create_pdf(topic)
            else:
                print("⚠️ Please provide a topic. Example: make pdf: Machine Learning")

        else:
            if should_search(user_input):
                print("🌐 Internet search triggered!")
            response = chat_with_search(user_input)
            print(f"\n🤖 AI: {response}\n")
            print("-" * 55)